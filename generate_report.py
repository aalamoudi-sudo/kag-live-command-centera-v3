#!/usr/bin/env python3
"""
generate_report.py
يولّد تقارير PPTX باستخدام python-pptx مباشرة على القوالب الأصلية
بدون أي تعديل على البنية — لا Repair، لا مشاكل
"""
import sys, json, copy, os
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
import io

TEMPLATES_DIR = os.path.join(os.path.dirname(__file__), "templates")

ORIGINALS = {
    "comprehensive": os.path.join(TEMPLATES_DIR, "comprehensive.pptx"),
    "أ":             os.path.join(TEMPLATES_DIR, "track-a.pptx"),
    "ب":             os.path.join(TEMPLATES_DIR, "track-b.pptx"),
    "ج":             os.path.join(TEMPLATES_DIR, "track-c.pptx"),
    "د":             os.path.join(TEMPLATES_DIR, "track-d.pptx"),
}

# ============================================================
# مساعدات
# ============================================================
def today_ar():
    return datetime.now().strftime("%Y-%m-%d")

def week_num():
    return datetime.now().isocalendar()[1]

def fmt_date(d):
    if not d: return ""
    try:
        return datetime.strptime(d[:10], "%Y-%m-%d").strftime("%Y-%m-%d")
    except:
        return str(d)

def status_label(s):
    if not s: return "—"
    if s in ["مكتملة","مكتمل","معتمدة","معتمد","ضمن المسار"]: return "أخضر ✓"
    if s in ["قيد التنفيذ","تحت المتابعة"]: return "أصفر"
    if s in ["معرضة للخطر","معرض للخطر","متأخرة"]: return "أحمر ✗"
    return s

def set_shape_text(shape, text):
    """استبدل نص shape مع الحفاظ على التنسيق الأصلي للـ run الأول"""
    if not shape.has_text_frame: return
    tf = shape.text_frame
    if not tf.paragraphs: return
    # احتفظ بالـ run الأول وتنسيقه
    para = tf.paragraphs[0]
    if not para.runs:
        # أضف run جديد
        from pptx.oxml.ns import qn
        from lxml import etree
        r = copy.deepcopy(para._p)
        return
    # امسح كل الـ paragraphs ما عدا الأول
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    # امسح كل الـ runs ما عدا الأول
    para = tf.paragraphs[0]
    for run in para.runs[1:]:
        run._r.getparent().remove(run._r)
    # عيّن النص
    if para.runs:
        para.runs[0].text = text
    else:
        para.text = text

def find_shape(slide, name):
    """ابحث عن shape باسمه"""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None

def fill_shape(slide, name, text):
    """ابحث عن shape وعيّن نصه"""
    shape = find_shape(slide, name)
    if shape:
        set_shape_text(shape, str(text) if text else "—")

def fill_shape_multiline(slide, name, lines):
    """عيّن نص متعدد الأسطر في shape"""
    shape = find_shape(slide, name)
    if not shape or not shape.has_text_frame: return
    tf = shape.text_frame
    # استخدم النص الأول كمرجع للتنسيق
    if not tf.paragraphs: return
    ref_para = tf.paragraphs[0]
    ref_run_fmt = ref_para.runs[0] if ref_para.runs else None

    # امسح كل المحتوى
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    para = tf.paragraphs[0]
    for run in para.runs[1:]:
        run._r.getparent().remove(run._r)

    text = "\n".join(lines) if lines else "—"
    if para.runs:
        para.runs[0].text = lines[0] if lines else "—"
    # أضف باقي الأسطر
    # للبساطة: ادمجها في سطر واحد مفصول بـ •
    combined = "  •  ".join(lines) if lines else "—"
    if para.runs:
        para.runs[0].text = combined

# ============================================================
# التقرير الشامل
# ============================================================
def fill_comprehensive(prs, state):
    tracks = state.get("tracks", [])
    items  = state.get("items",  [])
    risks  = [i for i in items if i.get("type") in ["مخاطرة","مخاطر","risks"] and i.get("status") != "مغلقة"]
    done   = [i for i in items if i.get("status") in ["مكتملة","مكتمل","معتمدة","معتمد"]]
    active = [i for i in items if i.get("status") == "قيد التنفيذ"]

    overall = round(sum(t.get("progress",0) for t in tracks) / len(tracks)) if tracks else 0
    worst = "أحمر" if any(t.get("status") in ["معرض للخطر","معرضة للخطر"] for t in tracks) \
        else "أصفر" if any(t.get("status") == "تحت المتابعة" for t in tracks) else "أخضر"

    def get_track(tid):
        return next((t for t in tracks if t.get("track") == tid), {})

    # شريحة 1 — الغلاف
    s1 = prs.slides[0]
    fill_shape(s1, "Text 4", f"الأسبوع {week_num()} ٢٠٢٦  |  الفترة: {today_ar()}")

    # شريحة 2 — الملخص التنفيذي
    s2 = prs.slides[1]
    fill_shape(s2, "Text 17", f"الحالة: {worst}  |  نسبة إنجاز اليوم: {overall}٪")
    done_items  = [i.get("title","") for i in done[:3]]
    active_items= [i.get("title","") for i in active[:3]]
    fill_shape(s2, "Text 21", "  •  ".join(done_items) if done_items else "لا يوجد")
    fill_shape(s2, "Text 25", "  •  ".join(active_items) if active_items else "لا يوجد")
    # القضايا
    for idx, r in enumerate(risks[:3]):
        base = 44 + idx * 28  # Text 44, 54, 64 ...
        fill_shape(s2, f"Text {base}", fmt_date(r.get("due","")))
        fill_shape(s2, f"Text {base+2}", r.get("owner","—"))
        fill_shape(s2, f"Text {base+4}", "متابعة عاجلة")
        fill_shape(s2, f"Text {base+6}", "تأثير على الجدول")
        fill_shape(s2, f"Text {base+8}", r.get("title","—"))

    # شريحة 3 — المسارات الأربعة
    s3 = prs.slides[2]
    track_shapes = {
        "أ": ("Text 32","Text 30","Text 28","Text 26","Text 24"),
        "ب": ("Text 44","Text 42","Text 40","Text 38","Text 36"),
        "ج": ("Text 56","Text 54","Text 52","Text 50","Text 48"),
        "د": ("Text 68","Text 66","Text 64","Text 62","Text 60"),
    }
    for tid, (s_amsl, s_alyom, s_ghad, s_hal, s_da3m) in track_shapes.items():
        t = get_track(tid)
        titems = [i for i in items if i.get("track") == tid]
        tdone  = [i for i in titems if i.get("status") in ["مكتملة","مكتمل","معتمدة","معتمد"]]
        tact   = [i for i in titems if i.get("status") == "قيد التنفيذ"]
        tnext  = sorted([i for i in titems if i.get("due") and i.get("due","") > today_ar()], key=lambda x: x.get("due",""))
        trisk  = [i for i in titems if i.get("type") in ["مخاطرة","مخاطر"] and i.get("status") != "مغلقة"]
        fill_shape(s3, s_amsl, tdone[0].get("title","—")  if tdone  else "—")
        fill_shape(s3, s_alyom,tact[0].get("title","—")   if tact   else "—")
        fill_shape(s3, s_ghad, tnext[0].get("title","—")  if tnext  else "—")
        fill_shape(s3, s_hal,  status_label(t.get("status","")))
        fill_shape(s3, s_da3m, f"{len(trisk)} مخاطر" if trisk else "لا يوجد")

    # شريحة 4 — السلامة
    s4 = prs.slides[3]
    crit = [r for r in risks if r.get("status") in ["معرضة للخطر","معرض للخطر"]]
    for idx, r in enumerate(crit[:2]):
        fill_shape(s4, f"Text {67 + idx*12}", r.get("title","—"))
        fill_shape(s4, f"Text {69 + idx*12}", r.get("owner","—"))

    # شريحة 5 — المخاطر والقرارات
    s5 = prs.slides[4]
    buckets = {
        "ح":  [r for r in risks if r.get("status") in ["معرضة للخطر","معرض للخطر"]],
        "أص": [r for r in risks if r.get("status") in ["تحت المتابعة","قيد التنفيذ"]],
        "أ":  [r for r in risks if r.get("status") not in ["معرضة للخطر","معرض للخطر","تحت المتابعة","قيد التنفيذ"]],
    }
    b_shapes = {"ح": ("Text 28","Text 22","Text 24","Text 26"), "أص": ("Text 38","Text 32","Text 34","Text 36"), "أ": ("Text 48","Text 42","Text 44","Text 46")}
    for key, (sq, sm, st, stow) in b_shapes.items():
        r = buckets[key][0] if buckets[key] else None
        fill_shape(s5, sq,   r.get("title","—")    if r else "—")
        fill_shape(s5, sm,   r.get("owner","—")    if r else "—")
        fill_shape(s5, st,   fmt_date(r.get("due","")) if r else "—")
        fill_shape(s5, stow, "متابعة عاجلة"         if r else "—")
    for idx, r in enumerate(risks[:3]):
        fill_shape(s5, f"Text {69 + idx*16}", r.get("title","—"))
        fill_shape(s5, f"Text {67 + idx*16}", r.get("title","—"))
        fill_shape(s5, f"Text {65 + idx*16}", "تأثير على الجدول")
        fill_shape(s5, f"Text {63 + idx*16}", "متابعة")

    # شريحة 6 — الجدول الزمني
    s6 = prs.slides[5]
    fill_timeline_slide(s6, items, None)

def fill_timeline_slide(slide, items, track_filter):
    filtered = [i for i in items if i.get("type") not in ["مخاطرة","مخاطر","risks"]]
    if track_filter:
        filtered = [i for i in filtered if i.get("track") == track_filter]
    done    = [i for i in filtered if i.get("status") in ["مكتملة","مكتمل","معتمدة","معتمد"]][:4]
    active  = [i for i in filtered if i.get("status") == "قيد التنفيذ"][:4]
    upcoming= sorted([i for i in filtered if i.get("due","") > today_ar()], key=lambda x:x.get("due",""))[:4]

    def fmt_list(lst): return "  •  ".join(i.get("title","") for i in lst) if lst else "لا يوجد"
    fill_shape(slide, "Text 42", fmt_list(done))
    fill_shape(slide, "Text 31", fmt_list(active))
    fill_shape(slide, "Text 20", fmt_list(upcoming))

# ============================================================
# تقرير المسار
# ============================================================
def fill_track(prs, track_key, state):
    tracks   = state.get("tracks", [])
    items    = state.get("items",  [])
    track    = next((t for t in tracks if t.get("track") == track_key), {})
    titems   = [i for i in items if i.get("track") == track_key]
    risks    = [i for i in titems if i.get("type") in ["مخاطرة","مخاطر","risks"] and i.get("status") != "مغلقة"]
    done     = [i for i in titems if i.get("status") in ["مكتملة","مكتمل","معتمدة","معتمد"]]
    active   = [i for i in titems if i.get("status") == "قيد التنفيذ"]
    upcoming = sorted([i for i in titems if i.get("due","") > today_ar()], key=lambda x:x.get("due",""))
    progress = track.get("progress", 0)
    status   = track.get("status", "")

    # شريحة 1 — الغلاف (لا تغيير)

    # شريحة 2 — الملخص
    s2 = prs.slides[1]
    fill_shape(s2, "Text 12", f"الحالة المختارة: {status_label(status)}  |  نسبة إنجاز اليوم: {progress}٪")
    fill_shape(s2, "Text 15",
               (done[0].get("title","") if done else "—") + "  |  " +
               (done[1].get("title","") if len(done)>1 else "—") + "  |  " +
               (done[2].get("title","") if len(done)>2 else "—"))
    fill_shape(s2, "Text 18",
               (risks[0].get("title","") if risks else "لا يوجد دعم عاجل") + "  |  " +
               "آخر موعد: " + fmt_date(risks[0].get("due","")) if risks else "لا يوجد دعم عاجل")

    # صفوف التحديث التفصيلي
    row_map = {
        "Text 42": done[0].get("title","—")    if done   else "—",
        "Text 40": active[0].get("title","—")  if active else "—",
        "Text 38": upcoming[0].get("title","—") if upcoming else "—",
        "Text 54": done[1].get("title","—")    if len(done)>1 else "—",
        "Text 52": active[1].get("title","—")  if len(active)>1 else "—",
        "Text 50": upcoming[1].get("title","—") if len(upcoming)>1 else "—",
    }
    for sname, val in row_map.items():
        fill_shape(s2, sname, val)

    # شريحة 3 — تفاصيل الأنشطة
    s3 = prs.slides[2]
    all_tasks = [i for i in titems if i.get("type") not in ["مخاطرة","مخاطر","risks"]][:8]
    shape_rows = [
        ("Text 35","Text 33","Text 31","Text 29","Text 27","Text 25"),
        ("Text 51","Text 49","Text 47","Text 45","Text 43","Text 41"),
        ("Text 67","Text 65","Text 63","Text 61","Text 59","Text 57"),
        ("Text 83","Text 81","Text 79","Text 77","Text 75","Text 73"),
        ("Text 99","Text 97","Text 95","Text 93","Text 91","Text 89"),
        ("Text 115","Text 113","Text 111","Text 109","Text 107","Text 105"),
        ("Text 131","Text 129","Text 127","Text 125","Text 123","Text 121"),
        ("Text 147","Text 145","Text 143","Text 141","Text 139","Text 137"),
    ]
    for idx, task in enumerate(all_tasks):
        if idx >= len(shape_rows): break
        st, sown, sstart, send, shal, spct = shape_rows[idx]
        fill_shape(s3, st,    task.get("title","—"))
        fill_shape(s3, sown,  task.get("owner","—"))
        fill_shape(s3, sstart,fmt_date(task.get("due","")))
        fill_shape(s3, send,  fmt_date(task.get("due","")))
        fill_shape(s3, shal,  status_label(task.get("status","")))
        fill_shape(s3, spct,  f"{task.get('progress',0)}٪")

    # شريحة 4 — المخاطر
    s4 = prs.slides[3]
    crit = [r for r in risks if r.get("status") in ["معرضة للخطر","معرض للخطر"]]
    med  = [r for r in risks if r.get("status") in ["تحت المتابعة","قيد التنفيذ"]]
    low  = [r for r in risks if r not in crit and r not in med]

    def fill_risk_row(slide, q_shape, own_shape, time_shape, tow_shape, r):
        if r:
            fill_shape(slide, q_shape,    r.get("title","—"))
            fill_shape(slide, own_shape,  r.get("owner","—"))
            fill_shape(slide, time_shape, fmt_date(r.get("due","")))
            fill_shape(slide, tow_shape,  "متابعة عاجلة")

    fill_risk_row(s4, "Text 25", "Text 21", "Text 23", "Text 23", crit[0] if crit else None)
    fill_risk_row(s4, "Text 37", "Text 33", "Text 35", "Text 35", med[0]  if med  else None)
    fill_risk_row(s4, "Text 49", "Text 45", "Text 47", "Text 47", low[0]  if low  else None)

    for idx, r in enumerate(risks[:3]):
        fill_shape(s4, f"Text {71 + idx*16}", r.get("title","—"))
        fill_shape(s4, f"Text {69 + idx*16}", r.get("title","—"))
        fill_shape(s4, f"Text {67 + idx*16}", "تأثير على الجدول")
        fill_shape(s4, f"Text {65 + idx*16}", "متابعة")

    # شريحة 5 — السلامة (نبقيها كما هي)

    # شريحة 6 — الجدول الزمني
    s6 = prs.slides[5]
    fill_timeline_track(s6, titems)

def fill_timeline_track(slide, items):
    filtered = [i for i in items if i.get("type") not in ["مخاطرة","مخاطر","risks"]]
    done    = [i for i in filtered if i.get("status") in ["مكتملة","مكتمل","معتمدة","معتمد"]][:4]
    active  = [i for i in filtered if i.get("status") == "قيد التنفيذ"][:4]
    upcoming= sorted([i for i in filtered if i.get("due","") > today_ar()], key=lambda x:x.get("due",""))[:4]

    def fmt(lst): return "  •  ".join(i.get("title","") for i in lst) if lst else "لا يوجد"
    fill_shape(slide, "Text 33", fmt(done))
    fill_shape(slide, "Text 24", fmt(active))
    fill_shape(slide, "Text 15", fmt(upcoming))

# ============================================================
# الدالة الرئيسية
# ============================================================
def generate_report(report_type, state):
    tpl_path = ORIGINALS.get(report_type)
    if not tpl_path or not os.path.exists(tpl_path):
        raise ValueError(f"قالب غير موجود: {report_type}")

    prs = Presentation(tpl_path)

    if report_type == "comprehensive":
        fill_comprehensive(prs, state)
    else:
        fill_track(prs, report_type, state)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

if __name__ == "__main__":
    # اقرأ النوع والبيانات من stdin
    data = json.loads(sys.stdin.read())
    report_type = data["type"]
    state       = data["state"]
    result      = generate_report(report_type, state)
    sys.stdout.buffer.write(result)
